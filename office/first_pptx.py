#!/usr/bin/env python
# -*- coding: utf-8 -*-
'''
@File    :   first_pptx.py
@Time    :   2024/04/12 08:56:39
@Author  :   hgh 
@Version :   1.0
@Desc    :    python 操作PPT, 官方链接文档: https://python-pptx.readthedocs.io/en/latest/
'''

# import module
import os
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt

import cv2 as cv

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from pptx.dml.fill import FillFormat
from pptx.enum.dml import MSO_FILL_TYPE, MSO_THEME_COLOR



def create_ppt_file(save_folder_path='./results', ppt_name = 'new_file.pptx'):
    
    if not os.path.exists(save_folder_path):
        os.mkdir(save_folder_path, mode=0o771)
        print(f'create {save_folder_path} success!')
        
    save_folder_path = os.path.join(save_folder_path, ppt_name)
    
    prs = Presentation()
    
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = 'this is my first ppt!'
    subtitle.text = 'sub title text was here~'
    
    # 获取slide对象
    slide_table = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 设置表格位置和大小
    left, top, width, height = Cm(7), Cm(6), Cm(12.8), Cm(5)
    
    # 设置表格行数、列数及大小
    shape_table = slide_table.shapes.add_table(7, 4, left, top, width, height)
    
    table = shape_table.table
    
    # 设置列宽度，行高
    table.columns[0].width = Cm
    
    
    prs.save(save_folder_path)


# ppt 插入图片+圆形图案
def insert_ppt_shape(r_value, save_folder_path='./results', ppt_name = 'new_file.pptx'):
    
    if not os.path.exists(save_folder_path):
        os.mkdir(save_folder_path, mode=0o771)
        print(f'create {save_folder_path} success!')
        
    save_folder_path = os.path.join(save_folder_path, ppt_name)
    
    prs = Presentation()
    
    title_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_slide_layout)

    # img_name = 'res.png'
    img_name = '01.png'
    x_step = Inches(2)  # 图片之间的水平间距
    y_step = Inches(2)  # 图片之间的垂直间距
    x = Inches(1)  # 图片左边缘起始位置
    y = Inches(1)  # 图片上边缘起始位置    
    slide.shapes.add_picture(os.path.join('.', img_name), x, y)
    
    shapes = slide.shapes
    
    left = Inches(1)
    top = Inches(1)
    width = Cm(r_value)
    height = Cm(r_value)
    
    circle = shapes.add_shape(
        autoshape_type_id=MSO_SHAPE.OVAL,   # 圆形
        left=left,
        top=top,
        width=width,
        height=height
    )
    
    # circle.text = 'ppt shape'
    # 设置填充为透明
    circle.fill.solid()
    # circle.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 使用ARGB格式，前三个为RGB颜色，最后一个值为透明度（0为完全透明）
    
    #设置透明色
    circle.fill.background()

    # 设置形状填充为透明（无填充）
    # circle.fill.type = MSO_FILL_TYPE
    
    # # Set the fill type to BACKGROUND for transparency
    # fill_format = FillFormat(shape.fill)
    # fill_format.type = MSO_FILL_TYPE.BACKGROUND
        
    # 设置边框颜色和线宽
    circle.line.color.rgb = RGBColor(0, 0, 0)  # 黑色边框，可替换为其他颜色
    circle.line.width = Pt(1)  # 边框宽度，单位为磅（Pt）
    
    # 如果需要，还可以设置边框样式，如 dashed, dotted, etc.
    # shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH  # 例如设置为虚线
    
    # shape.set_widget_bg(bg_rgb_color=[255, 255, 255])
    
    # txBox = shapes.add_textbox(left, top, width, height)
    # tf = txBox.text_frame

    prs.save(save_folder_path)


if __name__ == '__main__':
    
    
    # create_ppt_file()
    r_value = 20
    insert_ppt_shape(r_value, ppt_name='shape.pptx')
    
    
    
    pass





